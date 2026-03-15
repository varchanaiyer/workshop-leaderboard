"""
Build a polished PPTX from the workshop slide deck.
Uses python-pptx with custom styling for a professional look.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ─── BRAND COLORS ───────────────────────────────────────────────────────
BG_DARK = RGBColor(0x0F, 0x17, 0x2A)
BG_SECTION = RGBColor(0x1A, 0x1A, 0x2E)
ACCENT = RGBColor(0x6C, 0x63, 0xFF)
ACCENT2 = RGBColor(0x00, 0xD4, 0xAA)
TEXT_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_LIGHT = RGBColor(0xCC, 0xCC, 0xDD)
TEXT_DIM = RGBColor(0x88, 0x88, 0xAA)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
TABLE_HEADER = RGBColor(0x2D, 0x2B, 0x55)
TABLE_ROW1 = RGBColor(0x16, 0x16, 0x2A)
TABLE_ROW2 = RGBColor(0x1C, 0x1C, 0x32)
HIGHLIGHT = RGBColor(0xFF, 0xD7, 0x00)


def set_slide_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullets(slide, items, left, top, width, height,
                font_size=16, color=TEXT_LIGHT, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
    return txBox


def add_code_block(slide, code_text, left, top, width, height, font_size=11):
    shape = add_shape_bg(slide, left - Inches(0.1), top - Inches(0.1),
                         width + Inches(0.2), height + Inches(0.2), CODE_BG)
    shape.line.color.rgb = RGBColor(0x44, 0x44, 0x66)
    shape.line.width = Pt(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code_text
    p.font.size = Pt(font_size)
    p.font.color.rgb = RGBColor(0xA9, 0xDC, 0x76)
    p.font.name = "Courier New"
    return txBox


def add_table(slide, rows, left, top, width=None):
    n_rows = len(rows)
    n_cols = len(rows[0])
    w = width or Inches(12)
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, w, Inches(0.38 * n_rows))
    table = table_shape.table
    col_w = int(w / n_cols)
    for i in range(n_cols):
        table.columns[i].width = col_w
    for r_idx, row_data in enumerate(rows):
        for c_idx, cell_text in enumerate(row_data):
            cell = table.cell(r_idx, c_idx)
            cell.text = str(cell_text)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(11)
                paragraph.font.name = "Calibri"
                paragraph.font.bold = (r_idx == 0)
                paragraph.font.color.rgb = TEXT_WHITE if r_idx == 0 else TEXT_LIGHT
            cell.fill.solid()
            if r_idx == 0:
                cell.fill.fore_color.rgb = TABLE_HEADER
            elif r_idx % 2 == 1:
                cell.fill.fore_color.rgb = TABLE_ROW1
            else:
                cell.fill.fore_color.rgb = TABLE_ROW2
    return table_shape


# ─── SLIDE TEMPLATES ────────────────────────────────────────────────────

def make_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)
    add_text_box(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
                 "AI Research Implementation Workshop",
                 font_size=40, color=TEXT_WHITE, bold=True)
    add_text_box(slide, Inches(1), Inches(2.8), Inches(11), Inches(0.8),
                 "From Research Question to Paper \u2014 Using AI Coding Tools",
                 font_size=24, color=ACCENT2)
    add_text_box(slide, Inches(1), Inches(4.2), Inches(11), Inches(0.5),
                 "Algoverse AI Research  |  Single Session (~3 hrs)  |  20-30 Students",
                 font_size=16, color=TEXT_DIM)
    add_shape_bg(slide, Inches(0), Inches(7.0), Inches(13.33), Inches(0.06), ACCENT2)


def make_section_slide(prs, title, subtitle=""):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SECTION)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(0.12), Inches(7.5), ACCENT)
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11), Inches(1.2),
                 title, font_size=36, color=ACCENT2, bold=True)
    if subtitle:
        add_text_box(slide, Inches(1), Inches(3.8), Inches(11), Inches(0.8),
                     subtitle, font_size=18, color=TEXT_DIM)


def slide_base(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.04), ACCENT)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 title, font_size=28, color=TEXT_WHITE, bold=True)
    add_shape_bg(slide, Inches(0.6), Inches(1.05), Inches(4), Inches(0.03), ACCENT2)
    return slide


def make_task_slide(prs, task_num, title, duration, part_a_title, part_a_items,
                    part_b_title, part_b_items, share_note=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.9), ACCENT)
    add_text_box(slide, Inches(0.6), Inches(0.1), Inches(12), Inches(0.7),
                 f"TASK {task_num}: {title} ({duration})",
                 font_size=28, color=TEXT_WHITE, bold=True)
    # Part A
    add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.4),
                 part_a_title, font_size=20, color=ACCENT2, bold=True)
    add_bullets(slide, part_a_items,
                Inches(0.8), Inches(1.55), Inches(11.5), Inches(2.2),
                font_size=14, spacing=Pt(4))
    mid_y = Inches(3.8)
    add_shape_bg(slide, Inches(0.6), mid_y, Inches(12), Inches(0.02),
                 RGBColor(0x44, 0x44, 0x66))
    # Part B
    add_text_box(slide, Inches(0.6), mid_y + Inches(0.15), Inches(12), Inches(0.4),
                 part_b_title, font_size=20, color=ACCENT2, bold=True)
    add_bullets(slide, part_b_items,
                Inches(0.8), mid_y + Inches(0.6), Inches(11.5), Inches(2.2),
                font_size=14, spacing=Pt(4))
    if share_note:
        add_text_box(slide, Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
                     share_note, font_size=13, color=HIGHLIGHT, bold=True)
    return slide


WARNING_BG = RGBColor(0x2A, 0x15, 0x15)
WARNING_ACCENT = RGBColor(0xFF, 0x6B, 0x6B)


def make_mistakes_slide(prs, task_num, mistakes):
    """Create a 'Common Mistakes' slide with a table of mistakes and fixes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.04), WARNING_ACCENT)
    add_text_box(slide, Inches(0.6), Inches(0.3), Inches(12), Inches(0.7),
                 f"Common Mistakes \u2014 Task {task_num}",
                 font_size=28, color=WARNING_ACCENT, bold=True)
    add_shape_bg(slide, Inches(0.6), Inches(1.05), Inches(4), Inches(0.03), WARNING_ACCENT)
    add_table(slide, mistakes, Inches(0.3), Inches(1.3), Inches(12.7))
    return slide


# ─── BUILD ──────────────────────────────────────────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # 1: TITLE
    make_title_slide(prs)

    # 2: What You'll Build
    s = slide_base(prs, "What You'll Build Today")
    add_bullets(s, [
        "Set up and compare 3 AI coding tools (free to paid)",
        "Design a research benchmark AND run a prompting ablation study",
        "Implement experimental pipelines with AI assistance",
        "Analyze results and draft paper-ready outputs",
        "Learn to optimize code for Lambda GPUs",
        "",
        "You will all get DIFFERENT results \u2014 that's the point."
    ], Inches(0.6), Inches(1.3), Inches(12), Inches(3))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "By the end: two working research pipelines you can extend into real papers.",
                 font_size=13, color=HIGHLIGHT, bold=True)

    # 3: The 3 Tools
    s = slide_base(prs, "The 3 Tools We'll Use Today")
    add_table(s, [
        ["Tool", "Cost", "How It Works", "Best For"],
        ["Cursor + Gemini", "$20/mo (Cursor sub)", "IDE-based, inline completions, chat. Gemini is free.", "Fast iteration, visual editing"],
        ["Claude Code Router + Gemini", "$20/mo (CC sub)", "Terminal. Routes simple\u2192Gemini, complex\u2192Claude.", "Budget-conscious, most tasks"],
        ["Claude Code (direct)", "$20/mo + API usage", "Terminal. Claude Opus/Sonnet for everything.", "Complex logic, best quality"],
    ], Inches(0.6), Inches(1.3))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "Rule of thumb: 1-sentence task \u2192 Gemini. Paragraph-level task \u2192 Claude.",
                 font_size=13, color=HIGHLIGHT, bold=True)

    # 4: Tool Tier Strategy
    s = slide_base(prs, "Tool Tier Strategy \u2014 When to Use What")
    add_code_block(s,
        "FREE (Gemini)              $20/mo (Router/Cursor)      $20+ (Claude Direct)\n"
        "  |                            |                           |\n"
        "  |- Brainstorming             |- Multi-file edits         |- Complex architecture\n"
        "  |- Boilerplate               |- Debugging                |- Plan mode\n"
        "  |- Data formatting           |- Pipeline code            |- Worktree parallelism\n"
        "  |- Simple scripts            |- Analysis code            |- Paper-quality output",
        Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.8))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "Start free, level up only when you need to.",
                 font_size=13, color=HIGHLIGHT, bold=True)

    # 5: Setup Check
    s = slide_base(prs, "Setup Check (5 min)")
    add_bullets(s, [
        "Everyone should have ONE of these ready:",
        "  1. Cursor installed with Gemini model selected",
        "  2. Claude Code installed (npm install -g @anthropic-ai/claude-code)",
        "  3. Both (ideal)",
    ], Inches(0.6), Inches(1.3), Inches(12), Inches(2))
    add_code_block(s,
        '# Claude Code test:\nclaude "Say hello and tell me what model you are"\n\n'
        '# Cursor test:\n# Open a file, press Cmd+K, type "write hello world in python"',
        Inches(0.8), Inches(3.8), Inches(11.5), Inches(2))

    # 6: Today's Projects - BOTH
    s = slide_base(prs, "Today's Research Projects \u2014 You'll Do BOTH")
    add_text_box(s, Inches(0.6), Inches(1.3), Inches(5.5), Inches(0.4),
                 "Project A: BenchmarkLite", font_size=20, color=ACCENT2, bold=True)
    add_bullets(s, [
        '"Do LLMs give consistent answers when',
        '  you rephrase the same question?"',
        "",
        "\u2022 Design 20 questions + 3 paraphrases each",
        "\u2022 Test across 2-3 models",
        "\u2022 Measure consistency score",
    ], Inches(0.6), Inches(1.8), Inches(5.5), Inches(4))
    add_text_box(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4),
                 "Project B: Prompting Ablation", font_size=20, color=ACCENT2, bold=True)
    add_bullets(s, [
        '"Which prompting strategy works best',
        '  for [your chosen task]?"',
        "",
        "\u2022 Pick a task + dataset from our suggestions",
        "\u2022 Test 4 strategies (zero-shot \u2192 CoT+SC)",
        "\u2022 Compare accuracy across models",
    ], Inches(6.8), Inches(1.8), Inches(5.5), Inches(4))
    add_shape_bg(s, Inches(6.5), Inches(1.3), Inches(0.02), Inches(5),
                 RGBColor(0x44, 0x44, 0x66))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "Both are real, publishable research directions. You'll build BOTH today.",
                 font_size=14, color=HIGHLIGHT, bold=True)

    # 7: Skills Taught
    s = slide_base(prs, "What Both Projects Teach")
    add_table(s, [
        ["Research Skill", "Project A (Benchmark)", "Project B (Ablation)"],
        ["Experiment design", "Design evaluation categories", "Design prompting conditions"],
        ["Data curation", "Create question + paraphrase sets", "Pick dataset + curate examples"],
        ["API integration", "Call multiple LLM APIs", "Call multiple LLM APIs"],
        ["Statistical analysis", "Consistency metrics", "Accuracy comparisons"],
        ["Paper writing", "Tables, figures, claims", "Tables, figures, claims"],
    ], Inches(0.6), Inches(1.3))

    # 8: Example Questions for Benchmark
    s = slide_base(prs, "Example Questions for the Benchmark (Inspiration)")
    add_table(s, [
        ["Category", "Example Question", "Paraphrase 1", "Paraphrase 2"],
        ["Math", "What is 15% of 200?", "Calculate 15 percent of 200.", "200 \u00d7 0.15 = ?"],
        ["Logic", "If all roses are flowers and some fade quickly, must some roses?",
         "Every rose is a flower. Some flowers fade. Does it follow for roses?",
         "All roses \u2286 flowers. Some fade. Must some roses fade?"],
        ["Factual", "What is the capital of Australia?",
         "Which city serves as Australia's capital?", "Name Australia's capital city."],
        ["Ethics", "Is it ethical to use AI for hiring?",
         "Should companies use AI to decide who gets hired?",
         "Is there a moral issue with AI screening candidates?"],
        ["Creative", "Write a one-sentence story about a robot learning to paint.",
         "In one sentence, tell a story of a robot discovering painting.",
         "Compose a one-liner about an automaton taking up art."],
    ], Inches(0.3), Inches(1.3), Inches(12.7))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "Pick YOUR OWN categories \u2014 medical, legal, coding, history, psychology... The more varied, the better!",
                 font_size=13, color=HIGHLIGHT, bold=True)

    # 9: Example Ablation Walkthrough
    s = slide_base(prs, "Example Ablation \u2014 Same Question, 4 Strategies")
    add_text_box(s, Inches(0.6), Inches(1.2), Inches(12), Inches(0.4),
                 'Question: "A store sells a jacket for $120 after a 25% discount. Original price?"',
                 font_size=16, color=TEXT_LIGHT, bold=True)
    add_table(s, [
        ["Strategy", "What You Send", "Example Model Output"],
        ["Zero-shot", "Just the question", '"The original price was $160."'],
        ["Few-shot", "3 worked examples + question", '"Following the pattern: $120 / 0.75 = $160."'],
        ["Chain-of-Thought", 'Question + "Let\'s think step by step."',
         '"Step 1: Sale = 75% of original. Step 2: $120 = 0.75\u00d7X. Step 3: X = $160."'],
        ["CoT + Self-Consistency", "Run CoT 5 times, majority vote",
         "Runs: [$160, $160, $160, $150, $160] \u2192 Majority: $160"],
    ], Inches(0.3), Inches(1.8), Inches(12.7))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "CoT gives reasoning paths. Self-Consistency catches errors. Your results will vary!",
                 font_size=13, color=HIGHLIGHT, bold=True)

    # 10: Suggested Datasets
    s = slide_base(prs, "Suggested Datasets for the Ablation \u2014 Pick One")
    add_table(s, [
        ["Task", "Dataset", "Size", "How to Load", "Why It's Good"],
        ["Math word problems", "GSM8K", "50-100", 'load_dataset("gsm8k", "main", split="test[:50]")', "Clear answers, varying difficulty"],
        ["Sentiment", "SST-2", "50-100", 'load_dataset("glue", "sst2", split="validation[:50]")', "Binary labels, easy to score"],
        ["Reading comp.", "SQuAD 2.0", "50", 'load_dataset("squad_v2", split="validation[:50]")', "Extractive QA, gold answers"],
        ["Commonsense", "HellaSwag", "50", 'load_dataset("hellaswag", split="validation[:50]")', "Multiple choice, easy accuracy"],
        ["Code gen", "HumanEval", "20", 'load_dataset("openai_humaneval", split="test[:20]")', "Has test cases for verification"],
        ["Trivia", "TriviaQA", "50", 'load_dataset("trivia_qa", "rc", split="val[:50]")', "Short answers, exact match"],
    ], Inches(0.1), Inches(1.2), Inches(13))
    add_code_block(s,
        "# One line to load any dataset:\nfrom datasets import load_dataset\nds = load_dataset(\"gsm8k\", \"main\", split=\"test[:50]\")",
        Inches(0.8), Inches(5.5), Inches(8), Inches(1.2), font_size=12)

    # PART 1: PLAN
    make_section_slide(prs, "PART 1: Plan Your Research",
                       "Research planning with AI tools (30 min)")

    # 11: Planning Demo
    s = slide_base(prs, "Step 1 \u2014 Research Planning with AI")
    add_text_box(s, Inches(0.6), Inches(1.3), Inches(5.5), Inches(0.4),
                 "Cursor + Gemini (free)", font_size=18, color=ACCENT2, bold=True)
    add_bullets(s, [
        "Open research_plan.md, Cmd+K:",
        '"I want to evaluate LLM consistency on',
        "paraphrased questions. Help me design:",
        "(1) 5 question categories",
        "(2) how to create paraphrases",
        "(3) what metrics to measure",
        '(4) which models to test"',
    ], Inches(0.6), Inches(1.8), Inches(5.5), Inches(4.5), font_size=14)
    add_text_box(s, Inches(6.8), Inches(1.3), Inches(5.5), Inches(0.4),
                 "Claude Code (terminal)", font_size=18, color=ACCENT2, bold=True)
    add_bullets(s, [
        "In terminal:",
        'claude "I want to run an ablation study',
        "on prompting strategies for math using",
        "GSM8K. Use plan mode. Help me:",
        "(1) which strategies to compare",
        "(2) how to load the dataset",
        '(3) what metrics"',
        "",
        "Plan mode = structured output!",
    ], Inches(6.8), Inches(1.8), Inches(5.5), Inches(4.5), font_size=14)
    add_shape_bg(s, Inches(6.5), Inches(1.3), Inches(0.02), Inches(5),
                 RGBColor(0x44, 0x44, 0x66))

    # 12: CLAUDE.md
    s = slide_base(prs, "The CLAUDE.md \u2014 Your Research Lab Notebook")
    add_bullets(s, ["Before coding, create a CLAUDE.md in your project root:"],
                Inches(0.6), Inches(1.3), Inches(12), Inches(0.5))
    add_code_block(s,
        "# My Research Project\n\n"
        "## Research Question\n[Your 1-sentence question here]\n\n"
        "## Experiment Plan\n- Models: [list]\n- Dataset: [description]\n- Metrics: [list]\n\n"
        "## Decisions Made\n- [Date]: Chose X because Y\n\n"
        "## Mistakes to Avoid\n- [Claude adds these as you work]",
        Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 'Boris Cherny (Claude Code creator): "After every correction, update CLAUDE.md so that mistake doesn\'t repeat."',
                 font_size=13, color=HIGHLIGHT, bold=True)

    # 13: TASK 1
    make_task_slide(prs, 1, "Plan Both Experiments", "10 min",
        "Part A: Benchmark Plan (5 min)",
        [
            "1. Choose 5 question categories (see examples slide for inspiration)",
            "2. Write 1 example question + 3 paraphrases for one category",
            "3. Pick 2-3 models to test (e.g., GPT-4o, Claude Sonnet, Gemini)",
            "4. Define your consistency metric (exact match? semantic similarity?)",
            "Everyone's categories will differ \u2014 that's the point!",
        ],
        "Part B: Ablation Plan (5 min)",
        [
            "1. Choose your task + dataset (see datasets slide for options)",
            "2. Design 3-4 prompting strategies with example prompts",
            "3. Use the same 2-3 models from Part A",
            "4. Define your primary metric (accuracy, F1, etc.)",
            "Everyone's task/dataset choice will differ \u2014 great for comparison!",
        ],
        share_note="SHARE: 3-4 students present plans. Compare the variety!"
    )

    # Common Mistakes — Task 1
    make_mistakes_slide(prs, 1, [
        ["Mistake", "Example", "Fix"],
        ["All questions same difficulty",
         "5 easy math Qs like '2+3=?' \u2014 all get 100% consistency, learns nothing",
         "Mix: '2+3' AND 'integral of e^x sin(x)' in same category"],
        ["Paraphrases are just synonym swaps",
         "'What is 15% of 200?' \u2192 'What is fifteen percent of 200?' (trivial change)",
         "Restructure: 'If you take 15% from 200, what remains?' (different framing)"],
        ["No ground truth for subjective Qs",
         "'Is AI ethical?' \u2014 two valid but different essays score as 'inconsistent'",
         "Define per type: exact match for math, theme overlap for ethics"],
        ["Testing only 1 model",
         "GPT-4o gets 85% consistency \u2014 is that good or bad? No way to know.",
         "Compare GPT-4o (85%) vs Claude (72%) vs Gemini (90%) \u2192 now it means something"],
        ["Metric chosen after seeing results",
         "You pick 'semantic similarity' because it makes your numbers look better",
         "Pre-register: 'exact match for factual, cosine sim >0.8 for open-ended'"],
        ["No answer normalization",
         "'$160' vs 'The original price was $160.' flagged as mismatch (they agree!)",
         "Strip to key answer first: both become '$160', then compare"],
    ])

    # PART 2: IMPLEMENT
    make_section_slide(prs, "PART 2: Implement Your Pipelines",
                       "Building experimental code with AI (45 min)")

    # 14: Scaffolding
    s = slide_base(prs, "Step 2 \u2014 Project Scaffolding with Claude Code")
    add_code_block(s,
        'claude "Create a Python project for LLM evaluation with both\n'
        'a consistency benchmark and a prompting ablation. I need:\n'
        '- data/ for prompts and responses\n'
        '- src/ with: data loading, model API calls, metrics, viz\n'
        '- benchmark_main.py and ablation_main.py\n'
        '- requirements.txt\n'
        'Set up: python benchmark_main.py --model gpt-4 --dry-run"',
        Inches(0.8), Inches(1.5), Inches(11.5), Inches(3))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "Starter code is provided \u2014 but try scaffolding from scratch first!",
                 font_size=13, color=HIGHLIGHT, bold=True)

    # 15: Data Modules
    s = slide_base(prs, "The Data Modules")
    add_text_box(s, Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.4),
                 "Project A \u2014 Benchmark Dataset", font_size=18, color=ACCENT2, bold=True)
    add_code_block(s,
        '# data/questions.json\n'
        '{ "id": "math_001",\n'
        '  "category": "math",\n'
        '  "original": "What is 15% of 200?",\n'
        '  "paraphrases": [\n'
        '    "Calculate 15 percent of 200.",\n'
        '    "200 times 0.15 equals what?"\n'
        '  ] }',
        Inches(0.8), Inches(1.7), Inches(5.2), Inches(2.8), font_size=11)
    add_text_box(s, Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.4),
                 "Project B \u2014 Ablation Prompts + Dataset", font_size=18, color=ACCENT2, bold=True)
    add_code_block(s,
        '# src/prompts.py\n'
        'STRATEGIES = {\n'
        '  "zero_shot": "{question}",\n'
        '  "few_shot": "Ex:\\n{ex}\\nAnswer: {q}",\n'
        '  "cot": "{q}\\nThink step by step.",\n'
        '  "cot_sc": "run 5x, majority vote"\n'
        '}\n\n'
        'from datasets import load_dataset\n'
        'ds = load_dataset("gsm8k", split="test[:50]")',
        Inches(7.0), Inches(1.7), Inches(5.5), Inches(3.2), font_size=11)
    add_shape_bg(s, Inches(6.5), Inches(1.2), Inches(0.02), Inches(5),
                 RGBColor(0x44, 0x44, 0x66))

    # 16: Tool Comparison
    s = slide_base(prs, "Tool Comparison \u2014 Same Task, 3 Tools")
    add_table(s, [
        ["Task", "Cursor + Gemini", "Claude Router", "Claude Code Direct"],
        ["Generate data file", "Good (Cmd+K inline)", "Good (free via Gemini)", "Overkill"],
        ["Write API module", "Good", "Good", "Best (retry + caching)"],
        ["Multi-file pipeline", "Good", "Good", "Best (plan mode)"],
        ["Debug complex issue", "OK", "Route \u2192 Claude", "Best"],
        ["Quick iteration", "Best (IDE inline)", "Good", "Fine"],
    ], Inches(0.6), Inches(1.3))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "Demo: Watch us build the same API module in all 3 tools",
                 font_size=13, color=HIGHLIGHT, bold=True)

    # 17: TASK 2
    make_task_slide(prs, 2, "Implement Both Pipelines", "25 min",
        "Part A: Benchmark Pipeline (12 min)",
        [
            "1. Cursor/Gemini (5 min): Generate questions.json \u2014 20 questions with paraphrases",
            "2. Claude Code (5 min): Build evaluate.py \u2014 run all variants through model APIs",
            "3. Any tool (2 min): Write metrics.py \u2014 consistency scoring",
            "Test: python benchmark_main.py --dry-run --mock",
        ],
        "Part B: Ablation Pipeline (13 min)",
        [
            "1. Cursor/Gemini (3 min): Load dataset + create prompt templates",
            "2. Claude Code (7 min): Build ablation_evaluate.py \u2014 run all strategies",
            "3. Any tool (3 min): Add accuracy scoring + statistical comparison",
            "Test: python ablation_main.py --dry-run --mock",
        ],
        share_note="No API keys? Use --mock mode. Pipeline structure > real results."
    )

    # Common Mistakes — Task 2
    make_mistakes_slide(prs, 2, [
        ["Mistake", "Example", "Fix"],
        ["Not setting temperature/seed",
         "Run twice: GPT-4o gives 78% then 84% accuracy \u2014 which is real?",
         "Set temperature=0: now both runs give 81%, reproducible"],
        ["Few-shot examples from test set",
         "Your 3 few-shot examples are also in your 50 eval questions \u2192 inflated score",
         "Split: examples 1-5 for few-shot prompts, examples 6-55 for evaluation"],
        ["Different token limits per strategy",
         "CoT gets max_tokens=1000, zero-shot gets 100 \u2014 CoT wins because it can finish",
         "Set max_tokens=500 for ALL strategies, same playing field"],
        ["No rate limit handling",
         "Script crashes at question 23/50 with '429 Too Many Requests' \u2014 lose all progress",
         "Add: time.sleep(1) between calls + retry with backoff + save partial results"],
        ["No --mock mode for testing",
         "Debugging a JSON parsing bug costs $2 in API calls before you find the typo",
         "Build --mock first: return 'mock answer' \u2192 test full pipeline for $0"],
        ["Mixing model versions",
         "You tested 'gpt-4o' in Jan vs Mar \u2014 OpenAI updated it, results differ",
         "Log: model='gpt-4o-2024-11-20', timestamp='2025-03-01T10:00Z' in results"],
    ])

    # 18: Worktrees
    s = slide_base(prs, "The Worktree Power Move (Claude Code)")
    add_bullets(s, ["The #1 productivity unlock \u2014 run 3 parallel Claude sessions:"],
                Inches(0.6), Inches(1.3), Inches(12), Inches(0.5))
    add_code_block(s,
        "# Terminal 1 \u2014 benchmark pipeline\nclaude /worktree benchmark\n"
        '# "Build consistency benchmark: data, evaluation, metrics"\n\n'
        "# Terminal 2 \u2014 ablation pipeline\nclaude /worktree ablation\n"
        '# "Build prompting ablation: strategies, evaluation, accuracy"\n\n'
        "# Terminal 3 \u2014 shared API module\nclaude /worktree shared\n"
        '# "Build shared model API module with caching and retry"',
        Inches(0.8), Inches(2.0), Inches(11.5), Inches(3.5))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "3 parallel sessions, isolated branches. Merge when done.",
                 font_size=13, color=HIGHLIGHT, bold=True)

    # PART 3: ANALYZE
    make_section_slide(prs, "PART 3: Analyze & Write",
                       "From results to paper-ready outputs (30 min)")

    # 19: Analysis Demo
    s = slide_base(prs, "Step 3 \u2014 Results Analysis")
    add_code_block(s,
        'claude "Read results from both experiments.\n'
        'Benchmark: generate model x category consistency matrix.\n'
        'Ablation: generate strategy x model accuracy table.\n'
        'Create a bar chart for each. Run significance tests.\n'
        'Save figures to figures/ and summary to results/analysis.md"',
        Inches(0.8), Inches(1.5), Inches(11.5), Inches(2.2))
    add_bullets(s, [
        "",
        "Cursor alternative: open results JSON, highlight it, ask Gemini to",
        '"create a pandas analysis script that generates tables and plots."',
    ], Inches(0.6), Inches(4.2), Inches(12), Inches(1.5))

    # 20: Grill Me
    s = slide_base(prs, 'The "Grill Me" Pattern \u2014 AI as Reviewer')
    add_bullets(s, ["The most underused technique. Before writing your paper:"],
                Inches(0.6), Inches(1.3), Inches(12), Inches(0.5))
    add_code_block(s,
        'claude "Review my results in results/analysis.md. Be harsh.\n'
        '1. Do claims match the numbers in the tables?\n'
        '2. Statistical claims without significance tests?\n'
        '3. Anything misleading or cherry-picked?\n'
        '4. Biggest limitations I\'m not acknowledging?\n\n'
        'Grill me and don\'t let me write the paper until I pass."',
        Inches(0.8), Inches(2.0), Inches(11.5), Inches(2.8))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "Forces critical thinking \u2014 exactly what a real reviewer will do.",
                 font_size=13, color=HIGHLIGHT, bold=True)

    # 21: TASK 3
    make_task_slide(prs, 3, "Analyze Both & Compare", "15 min",
        "Part A: Benchmark Analysis (7 min)",
        [
            "1. Generate consistency matrix (Model \u00d7 Category)",
            "2. Create a grouped bar chart",
            "3. Write 2-sentence finding: What did you find? Was it expected?",
        ],
        "Part B: Ablation Analysis (8 min)",
        [
            "1. Generate accuracy table (Strategy \u00d7 Model)",
            "2. Create a grouped bar chart",
            "3. Write 2-sentence finding: Which strategy won? By how much?",
            "",
            "Combined: Do benchmark + ablation results tell a related story?",
        ],
        share_note="SHARE: 3-4 students present. Did anyone get opposite results? That's research!"
    )

    # Common Mistakes — Task 3
    make_mistakes_slide(prs, 3, [
        ["Mistake", "Example", "Fix"],
        ["No error bars",
         "Report 'CoT accuracy: 78%' \u2014 but across 50 Qs, 95% CI is [65%, 91%]",
         "Always: '78% (\u00b113%, n=50)' \u2014 shows the uncertainty honestly"],
        ["Cherry-picking results",
         "CoT wins on math (85%) but loses on sentiment (40%) \u2014 you only show math",
         "Report ALL: 'CoT wins on math (+15%) but hurts sentiment (-10%)'"],
        ["Claims don't match numbers",
         "'CoT significantly outperforms zero-shot' but scores are 78% vs 76% (p=0.42)",
         "Grill Me check: 'Is 2% with p=0.42 significant? No. Rephrase the claim.'"],
        ["Y-axis manipulation",
         "Bar chart Y-axis from 75% to 85% \u2014 makes 3% gap look like 10x difference",
         "Start Y-axis at 0, or add a note: 'Note: Y-axis starts at 75%'"],
        ["No combined insight",
         "Two separate findings but no connection between benchmark + ablation results",
         "'Models inconsistent on ethics also gain most from CoT on ethics' \u2192 real insight"],
        ["Overgeneralizing",
         "'LLMs are unreliable' from testing 2 models on 50 math questions",
         "'GPT-4o and Claude show 72-85% consistency on math paraphrases (n=50)'"],
    ])

    # PART 4: LAMBDA GPU
    make_section_slide(prs, "PART 4: Lambda GPU Optimization",
                       "Making your code fast and cost-efficient (25 min)")

    # 22: Why Lambda
    s = slide_base(prs, "Why Lambda GPUs?")
    add_bullets(s, [
        "Research beyond API calls needs GPUs:",
        "  \u2022 Fine-tuning models (LoRA, full fine-tune)",
        "  \u2022 Running local open-source models (Llama, Mistral)",
        "  \u2022 Large-scale inference for benchmarks",
        "  \u2022 Training from scratch",
        "",
        "Lambda gives you cloud GPUs on demand.",
        "But GPU time = money. Every wasted minute costs real dollars.",
    ], Inches(0.6), Inches(1.3), Inches(12), Inches(4))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 "Rule #1: Always nvidia-smi first. 0% GPU during training = something is wrong.",
                 font_size=13, color=HIGHLIGHT, bold=True)

    # 23-25: Optimizations 1-3
    for title, code, note in [
        ("Optimization 1: Don't Waste GPU on Data Loading",
         "# BAD \u2014 GPU sits idle waiting for data\n"
         "train_loader = DataLoader(dataset, batch_size=32)\n\n"
         "# GOOD \u2014 parallel data loading, GPU never waits\n"
         "train_loader = DataLoader(\n"
         "    dataset, batch_size=32,\n"
         "    num_workers=4,          # parallel CPU workers\n"
         "    pin_memory=True,        # faster CPU->GPU transfer\n"
         "    prefetch_factor=2       # pre-load next batches\n"
         ")",
         "Impact: 2-5x speedup for data-heavy workloads."),
        ("Optimization 2: Mixed Precision Training",
         "# BAD \u2014 full FP32, uses 2x memory, slower\n"
         "loss = model(batch); loss.backward()\n\n"
         "# GOOD \u2014 mixed precision, half memory, 2-3x faster\n"
         "from torch.amp import autocast, GradScaler\n\n"
         "scaler = GradScaler()\n"
         "with autocast(device_type='cuda'):    # FP16 forward\n"
         "    loss = model(batch)\n"
         "scaler.scale(loss).backward()          # scaled backward\n"
         "scaler.step(optimizer); scaler.update()",
         "Impact: 2x larger batches, 2-3x faster on A100/H100."),
        ("Optimization 3: Batch Your Inference",
         "# BAD \u2014 one prompt at a time (GPU 5% utilized)\n"
         "for prompt in prompts:\n"
         "    result = model.generate(prompt)\n\n"
         "# GOOD \u2014 batch inference (GPU 90%+ utilized)\n"
         "prompt_loader = DataLoader(prompts, batch_size=16)\n"
         "for batch in prompt_loader:\n"
         "    tok = tokenizer(batch, padding=True, return_tensors='pt').to('cuda')\n"
         "    outputs = model.generate(**tok)",
         "Impact: 10-20x faster for inference/evaluation workloads."),
    ]:
        s = slide_base(prs, title)
        add_code_block(s, code, Inches(0.8), Inches(1.5), Inches(11.5), Inches(3.5))
        add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                     note, font_size=13, color=HIGHLIGHT, bold=True)

    # 26: Monitor + Checkpoint
    s = slide_base(prs, "Optimization 4: Monitor and Checkpoint")
    add_text_box(s, Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.4),
                 "GPU Monitoring", font_size=18, color=ACCENT2, bold=True)
    add_bullets(s, [
        "Always watch GPU utilization:",
        "  watch -n 1 nvidia-smi",
        "", "If util < 50%: STOP and fix.",
        "  \u2022 Increase batch size",
        "  \u2022 Add num_workers",
        "  \u2022 Check CPU bottlenecks",
    ], Inches(0.6), Inches(1.7), Inches(5.5), Inches(4.5), font_size=14)
    add_text_box(s, Inches(6.8), Inches(1.2), Inches(5.5), Inches(0.4),
                 "Checkpointing", font_size=18, color=ACCENT2, bold=True)
    add_code_block(s,
        '# Save every epoch\ntorch.save({\n'
        '    "epoch": epoch,\n'
        '    "model": model.state_dict(),\n'
        '    "optim": optimizer.state_dict(),\n'
        '}, f"ckpt_{epoch}.pt")\n\n'
        '# Resume:\nckpt = torch.load("ckpt_5.pt")\n'
        'model.load_state_dict(ckpt["model"])',
        Inches(7.0), Inches(1.7), Inches(5.5), Inches(3.5), font_size=11)
    add_shape_bg(s, Inches(6.5), Inches(1.2), Inches(0.02), Inches(5),
                 RGBColor(0x44, 0x44, 0x66))

    # 27: Cost Sheet
    s = slide_base(prs, "Lambda Cost Cheat Sheet")
    add_table(s, [
        ["GPU", "$/hr (on-demand)", "$/hr (spot)", "Best For"],
        ["A10 (24GB)", "~$0.75", "~$0.30", "Small inference, fine-tuning \u22647B"],
        ["A100 (80GB)", "~$1.50", "~$0.60", "Training 7B-13B, large batch"],
        ["H100 (80GB)", "~$2.50", "~$1.00", "Training 13B+, fastest"],
    ], Inches(0.6), Inches(1.3))
    add_bullets(s, [
        "", "Budget rules:",
        "  1. Develop locally first. GPU only when code is tested.",
        "  2. Spot instances (3-5x cheaper) + checkpointing.",
        "  3. nvidia-smi every time \u2014 util < 50% = stop and optimize.",
        "  4. Set a Lambda budget alert.",
    ], Inches(0.6), Inches(3.8), Inches(12), Inches(2.5))

    # 28: TASK 4
    s = slide_base(prs, "TASK 4: Fix the Slow Training Script (10 min)")
    add_bullets(s, [
        "slow_train.py has at least 7 performance problems. Find them all.",
    ], Inches(0.6), Inches(1.3), Inches(12), Inches(0.5))
    add_code_block(s,
        "# Processes ONE sample at a time \u2014 find all issues:\n"
        "for i in range(len(texts)):\n"
        "    inputs = tokenizer(texts[i], return_tensors='pt').to('cuda')\n"
        "    outputs = model(**inputs, labels=torch.tensor([labels[i]]).to('cuda'))\n"
        "    outputs.loss.backward()\n"
        "    optimizer.step()\n"
        "    optimizer.zero_grad()",
        Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.5))
    add_text_box(s, Inches(0.6), Inches(6.5), Inches(12), Inches(0.5),
                 'Prompt: "This script has at least 7 performance problems. Find and fix them all."',
                 font_size=13, color=HIGHLIGHT, bold=True)

    # Common Mistakes — Task 4
    make_mistakes_slide(prs, 4, [
        ["Mistake", "Example", "Fix"],
        ["Batch size = 1",
         "for i in range(len(texts)): model(texts[i]) \u2014 nvidia-smi shows GPU 5%",
         "DataLoader(dataset, batch_size=32) \u2014 GPU jumps to 85%+"],
        ["num_workers=0",
         "DataLoader(ds, batch_size=32) \u2014 GPU waits 0.5s per batch for CPU to load data",
         "DataLoader(ds, batch_size=32, num_workers=4, pin_memory=True) \u2014 no wait"],
        ["No mixed precision",
         "A100 runs FP32: 8GB model, batch=16, 45 min/epoch",
         "Add autocast + GradScaler: 4GB model, batch=32, 15 min/epoch (3x faster)"],
        [".item()/.cpu() in loop",
         "loss_sum += loss.item()  # every step forces GPU\u2192CPU sync, kills pipeline",
         "loss_sum += loss.detach()  # log every 100 steps: if step % 100: print(loss_sum.item())"],
        ["Forgetting model.eval()",
         "Inference with dropout active: accuracy varies 60-75% across runs (should be 72%)",
         "model.eval() + with torch.no_grad(): \u2192 stable 72%, also 30% faster"],
        ["No checkpointing",
         "Spot instance dies at epoch 8/10 after 6 hours ($9 wasted), start over",
         "torch.save(ckpt, f'ckpt_{epoch}.pt') every epoch \u2192 resume from epoch 8"],
    ])

    # PART 5: WRAP-UP
    make_section_slide(prs, "PART 5: Wrap-Up & Next Steps",
                       "Your toolkit going forward (10 min)")

    # 29: Tool Matrix
    s = slide_base(prs, "Tool Decision Matrix \u2014 Your Cheat Sheet")
    add_table(s, [
        ["Task", "Cursor + Gemini", "Claude Router", "Claude Code"],
        ["Brainstorm", "Good", "Good (free)", "Overkill"],
        ["Boilerplate code", "Great (inline)", "Good (free)", "Overkill"],
        ["Experiment plan", "OK", "Good", "Best (plan mode!)"],
        ["Multi-file pipeline", "Good", "Good", "Best"],
        ["Debug complex issue", "OK", "Route \u2192 Claude", "Best"],
        ["Data formatting", "Best (IDE)", "Good (free)", "Fine"],
        ["Figures", "Good", "Good", "Good"],
        ["Review paper", "OK", "Route \u2192 Claude", "Best (grill!)"],
        ["Parallel work", "N/A", "Worktrees!", "Worktrees!"],
    ], Inches(0.6), Inches(1.3))

    # 30: Habits
    s = slide_base(prs, "CLAUDE.md Habits Going Forward")
    add_bullets(s, [
        "For every research project:",
        "",
        "1. Start with CLAUDE.md \u2014 define question, plan, and metrics",
        "2. Update after every correction \u2014 build institutional memory",
        "3. Add verification steps \u2014 tests that confirm output format",
        "4. Use plan mode for architecture \u2014 don't code without a plan",
        '5. Grill yourself \u2014 "Review my claims against my data. Be harsh."',
        "",
        "This separates students who use AI as a crutch",
        "from researchers who use AI as a force multiplier.",
    ], Inches(0.6), Inches(1.3), Inches(12), Inches(4.5))

    # 31: Next Steps
    s = slide_base(prs, "Next Steps for Your Algoverse Research")
    add_bullets(s, [
        "Today: You built TWO experiment pipelines in ~3 hours",
        "This week: Expand them \u2014 more questions, more models, more analysis",
        "Before your deadline: Use the 'grill me' pattern to self-review",
        "When submitting: Your CLAUDE.md becomes your reproducibility notes",
        "",
        "Upcoming deadlines:",
        "  \u2022 ACL 2026 workshops (March\u2013April deadlines)",
        "  \u2022 ICML 2026",
        "  \u2022 Check Conference Submissions tracker for your team",
    ], Inches(0.6), Inches(1.3), Inches(12), Inches(4.5))

    # 32: Resources
    s = slide_base(prs, "Resources")
    add_table(s, [
        ["Resource", "Link"],
        ["Claude Code docs", "docs.anthropic.com/en/docs/claude-code"],
        ["Claude Code tips (Boris Cherny)", "x.com/bcherny/status/2017742741636321619"],
        ["Cursor docs", "docs.cursor.com"],
        ["Lambda GPU guide", "docs.lambdalabs.com"],
        ["Algoverse Conference Tracker", "(internal Airtable)"],
    ], Inches(0.6), Inches(1.3))

    # 33: Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)
    add_shape_bg(slide, Inches(0), Inches(0), Inches(13.33), Inches(0.06), ACCENT)
    add_text_box(slide, Inches(1), Inches(2.0), Inches(11), Inches(1.2),
                 "Now Go Build Something.", font_size=44, color=TEXT_WHITE,
                 bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(3.5), Inches(11), Inches(0.8),
                 "Two research pipelines ready. Expand them, run real experiments, publish.",
                 font_size=20, color=ACCENT2, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
                 "Algoverse AI Research", font_size=16, color=TEXT_DIM,
                 alignment=PP_ALIGN.CENTER)
    add_shape_bg(slide, Inches(0), Inches(7.0), Inches(13.33), Inches(0.06), ACCENT2)

    output = "/Users/archana/Desktop/Algoverse/ai-research-workshop/slides/AI_Research_Workshop.pptx"
    prs.save(output)
    print(f"Saved to {output}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    build_presentation()
